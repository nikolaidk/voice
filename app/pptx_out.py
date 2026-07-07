"""Export a slide deck as a PowerPoint (.pptx) file.

Same theme as the web slideshow and video; each slide carries the narration
it accompanies as speaker notes, so the file works as a real presentation.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .script_gen import PodcastScript
from .slides import DEFAULT_THEME, SlideDeck, SlideTheme, slide_visual_path

_W_IN, _H_IN = 13.333, 7.5


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _font_name(css_stack: str) -> str:
    first = css_stack.split(",")[0].strip().strip("'\"")
    return first if first and not first.startswith("-") else "Helvetica Neue"


def build_pptx(
    title: str,
    deck: SlideDeck,
    script: PodcastScript,
    theme: SlideTheme | None,
    assets_dir: Path,
    out_path: Path,
    footer: str | None = None,
) -> None:
    theme = theme or DEFAULT_THEME
    panel = _rgb(theme.panel)
    ink = _rgb(theme.text_color)
    muted = _rgb(theme.muted_color)
    accent = _rgb(theme.accent)
    body_font = _font_name(theme.font_family)
    head_font = _font_name(theme.heading_font_family)

    prs = Presentation()
    prs.slide_width = Inches(_W_IN)
    prs.slide_height = Inches(_H_IN)
    blank = prs.slide_layouts[6]

    n = len(deck.slides)
    for i, slide in enumerate(deck.slides):
        s = prs.slides.add_slide(blank)

        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = panel

        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Emu(91440), Inches(_H_IN)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False

        is_cover = i == 0
        visual = None
        vp = slide_visual_path(slide, i, assets_dir, raster=True)
        if vp is not None:
            visual = vp

        text_w = Inches(6.1 if visual else 11.4)

        # title
        tb = s.shapes.add_textbox(Inches(0.9), Inches(0.55 if is_cover else 0.5),
                                  text_w, Inches(2.2 if is_cover else 1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide.title
        p.font.size = Pt(40 if is_cover else 30)
        p.font.bold = True
        p.font.name = head_font
        p.font.color.rgb = ink

        y = 2.6 if is_cover else 1.9

        if slide.big_statement:
            st = s.shapes.add_textbox(Inches(0.9), Inches(y), text_w, Inches(1.8))
            st_tf = st.text_frame
            st_tf.word_wrap = True
            sp = st_tf.paragraphs[0]
            sp.text = slide.big_statement
            sp.font.size = Pt(22 if is_cover else 20)
            sp.font.italic = True
            sp.font.name = body_font
            sp.font.color.rgb = accent if is_cover else muted
            y += 1.9

        if slide.bullets:
            dense = len(slide.bullets) > 5 or sum(len(b) for b in slide.bullets) > 260
            bt = s.shapes.add_textbox(Inches(0.9), Inches(y), text_w,
                                      Inches(_H_IN - y - 0.6))
            bt_tf = bt.text_frame
            bt_tf.word_wrap = True
            for k, bullet in enumerate(slide.bullets):
                bp = bt_tf.paragraphs[0] if k == 0 else bt_tf.add_paragraph()
                bp.text = f"•  {bullet}"
                bp.font.size = Pt(14 if dense else 18)
                bp.font.name = body_font
                bp.font.color.rgb = ink
                bp.space_after = Pt(8 if dense else 12)

        if visual is not None:
            from PIL import Image

            with Image.open(visual) as img:
                iw, ih = img.size
            box_x, box_y, box_w, box_h = 7.3, 1.0, 5.4, 5.5
            ratio = min(box_w / iw, box_h / ih)
            w_in, h_in = iw * ratio, ih * ratio
            s.shapes.add_picture(
                str(visual),
                Inches(box_x + (box_w - w_in) / 2),
                Inches(box_y + (box_h - h_in) / 2),
                Inches(w_in), Inches(h_in),
            )

        if footer:
            fb = s.shapes.add_textbox(Inches(0.9), Inches(_H_IN - 0.55),
                                      Inches(5.5), Inches(0.4))
            fp_ = fb.text_frame.paragraphs[0]
            fp_.text = footer
            fp_.font.size = Pt(10)
            fp_.font.name = body_font
            fp_.font.color.rgb = muted

        num = s.shapes.add_textbox(Inches(_W_IN - 1.5), Inches(_H_IN - 0.55),
                                   Inches(1.1), Inches(0.4))
        np_ = num.text_frame.paragraphs[0]
        np_.text = f"{i + 1} / {n}"
        np_.alignment = PP_ALIGN.RIGHT
        np_.font.size = Pt(12)
        np_.font.name = body_font
        np_.font.color.rgb = muted

        # narration this slide accompanies -> speaker notes
        start = slide.first_line
        end = deck.slides[i + 1].first_line if i + 1 < n else len(script.lines)
        notes = "\n\n".join(
            line.text for line in script.lines[start:end] if line.text.strip()
        )
        if notes:
            s.notes_slide.notes_text_frame.text = notes

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
